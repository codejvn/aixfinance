from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
BLACK    = RGBColor(0x0D, 0x0D, 0x0D)
SURFACE  = RGBColor(0x1A, 0x1A, 0x1A)
BRG      = RGBColor(0x2C, 0x5F, 0x2E)
BRG_LT   = RGBColor(0x3D, 0x7A, 0x40)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x9C, 0xA3, 0xAF)
GREEN_LT = RGBColor(0xD1, 0xFA, 0xD1)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = color
    return s


def text(slide, txt, left, top, width, height,
         size=20, bold=False, italic=False,
         color=WHITE, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = txt
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color


def bullets(slide, items, left, top, width, height, size=17, color=WHITE):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"\u2022  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(5)


def green_bar(slide):
    rect(slide, 0, Inches(0.52), W, Inches(0.04), BRG)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, BLACK)
rect(s, 0, 0, Inches(0.35), H, BRG)
rect(s, Inches(0.35), Inches(2.65), W, Inches(0.04), BRG)

text(s, "KAPPA",
     Inches(0.65), Inches(1.1), Inches(11), Inches(1.5),
     size=80, bold=True)
text(s, "Actor-Critic Multi-Agent Trading Portfolio",
     Inches(0.65), Inches(2.8), Inches(11), Inches(0.65),
     size=26, color=GREEN_LT)
rect(s, Inches(0.35), Inches(3.55), W, Inches(0.04), BRG)
text(s, "Perpetuals  \u00b7  Crypto  \u00b7  Real World Assets",
     Inches(0.65), Inches(3.75), Inches(11), Inches(0.55),
     size=18, italic=True, color=MUTED)
text(s, "Cross-Model Adversarial Ensemble  |  Claude \u00d7 GPT",
     Inches(0.65), Inches(6.65), Inches(11), Inches(0.55),
     size=13, color=MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — What is Kappa?
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, BLACK)
green_bar(s)
text(s, "What is Kappa?", Inches(0.6), Inches(0.8), Inches(11), Inches(0.7), size=36, bold=True)
rect(s, Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.03), BRG)

text(s,
     "Kappa is a consumer-facing web application that uses a Cross-Model Adversarial Ensemble "
     "to generate risk-adjusted thematic trading portfolios for perpetuals, crypto, and "
     "real-world assets.",
     Inches(0.6), Inches(1.8), Inches(12.1), Inches(1.2), size=19)

bullets(s, [
    "User defines a theme: AI, DeFi, Layer-1s, Geopolitics, or any custom topic",
    "AI agents research the market in real time using live news and price data",
    "Two AI models (Claude + GPT) compete to produce the optimal portfolio",
    "A Risk Analyst synthesises the best ideas into a final executable strategy",
    "The strategy is back-tested and presented with full performance metrics",
], Inches(0.6), Inches(3.1), Inches(12.1), Inches(4.0), size=19)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — The Problem
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, BLACK)
green_bar(s)
text(s, "The Problem", Inches(0.6), Inches(0.8), Inches(11), Inches(0.7), size=36, bold=True)
rect(s, Inches(0.6), Inches(1.6), Inches(3.5), Inches(0.03), BRG)

rect(s, Inches(6.6), Inches(1.75), Inches(0.03), Inches(4.8), BRG)

text(s, "For Retail Investors", Inches(0.6), Inches(1.75), Inches(5.8), Inches(0.5),
     size=20, bold=True, color=GREEN_LT)
bullets(s, [
    "No access to institutional-grade research tools",
    "Overwhelmed by market noise and conflicting signals",
    "Risk management is manual and error-prone",
    "No systematic way to translate a thesis into a portfolio",
], Inches(0.6), Inches(2.35), Inches(5.8), Inches(3.2), size=18)

text(s, "The Old Way", Inches(6.9), Inches(1.75), Inches(5.8), Inches(0.5),
     size=20, bold=True, color=GREEN_LT)
bullets(s, [
    "Single AI model = single point of bias",
    "No adversarial stress-testing of ideas",
    "Static portfolios with no back-test validation",
    "No sentiment + price data synthesis",
], Inches(6.9), Inches(2.35), Inches(5.8), Inches(3.2), size=18)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Architecture Pipeline
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, BLACK)
green_bar(s)
text(s, "Architecture Pipeline", Inches(0.6), Inches(0.8), Inches(11), Inches(0.7),
     size=36, bold=True)
rect(s, Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.03), BRG)

steps = [
    ("1", "User Input",          "Theme + risk tolerance entered via the Streamlit UI"),
    ("2", "Data Gathering",      "Live news via Tavily & NewsAPI  \u00b7  Price/volume via yfinance"),
    ("3", "Sentiment Synthesis", "NLP sentiment scored per article \u2192 unified JSON context block"),
    ("4", "Actor Arena",         "Claude Actor & GPT Actor each propose independent trade sets"),
    ("5", "Critic Arena",        "Claude Critic attacks GPT  \u00b7  GPT Critic attacks Claude"),
    ("6", "Risk Analyst",        "Final LLM synthesises the best logic \u2192 execution JSON"),
    ("7", "Backtesting",         "Portfolio simulated against historical data \u2192 equity curve & metrics"),
]

y = Inches(1.85)
for num, title, desc in steps:
    rect(s, Inches(0.6), y, Inches(0.44), Inches(0.44), BRG)
    text(s, num, Inches(0.6), y, Inches(0.44), Inches(0.44),
         size=13, bold=True, align=PP_ALIGN.CENTER)
    text(s, title, Inches(1.15), y, Inches(2.6), Inches(0.44),
         size=14, bold=True, color=GREEN_LT)
    text(s, desc, Inches(3.85), y, Inches(9.0), Inches(0.44), size=14, color=MUTED)
    y += Inches(0.72)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — The Actor-Critic Arena
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, BLACK)
green_bar(s)
text(s, "The Actor-Critic Arena", Inches(0.6), Inches(0.8), Inches(11), Inches(0.7),
     size=36, bold=True)
rect(s, Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.03), BRG)

# Claude box
rect(s, Inches(0.5), Inches(1.85), Inches(5.5), Inches(2.35), SURFACE)
rect(s, Inches(0.5), Inches(1.85), Inches(5.5), Inches(0.07), BRG)
text(s, "Claude (Anthropic)", Inches(0.7), Inches(2.0), Inches(5.1), Inches(0.5),
     size=17, bold=True, color=GREEN_LT)
bullets(s, [
    "Actor: proposes Long/Short trades with leverage",
    "Critic: stress-tests GPT portfolio for risk flaws",
], Inches(0.7), Inches(2.55), Inches(5.1), Inches(1.4), size=16)

text(s, "VS", Inches(6.1), Inches(2.5), Inches(1.15), Inches(0.7),
     size=28, bold=True, color=BRG, align=PP_ALIGN.CENTER)

# GPT box
rect(s, Inches(7.2), Inches(1.85), Inches(5.5), Inches(2.35), SURFACE)
rect(s, Inches(7.2), Inches(1.85), Inches(5.5), Inches(0.07), BRG)
text(s, "GPT (via Dedalus Labs)", Inches(7.4), Inches(2.0), Inches(5.1), Inches(0.5),
     size=17, bold=True, color=GREEN_LT)
bullets(s, [
    "Actor: proposes competing trade set independently",
    "Critic: stress-tests Claude portfolio for risk flaws",
], Inches(7.4), Inches(2.55), Inches(5.1), Inches(1.4), size=16)

# Arrows down to Risk Analyst
text(s, "\u2193", Inches(3.1), Inches(4.2), Inches(1.0), Inches(0.45),
     size=24, color=BRG, align=PP_ALIGN.CENTER)
text(s, "\u2193", Inches(9.55), Inches(4.2), Inches(1.0), Inches(0.45),
     size=24, color=BRG, align=PP_ALIGN.CENTER)

# Risk Analyst box
rect(s, Inches(2.5), Inches(4.65), Inches(8.3), Inches(2.0), SURFACE)
rect(s, Inches(2.5), Inches(4.65), Inches(8.3), Inches(0.07), BRG)
text(s, "Risk Analyst (Final Layer)", Inches(2.7), Inches(4.8), Inches(8.0), Inches(0.5),
     size=17, bold=True, color=GREEN_LT)
text(s,
     "Reviews both proposals + critiques \u2192 synthesises best logic \u2192 "
     "outputs execution JSON (Asset, Direction, Leverage, Stop-Loss, Allocation %)",
     Inches(2.7), Inches(5.35), Inches(7.8), Inches(1.1), size=15)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Data Sources
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, BLACK)
green_bar(s)
text(s, "Data Sources", Inches(0.6), Inches(0.8), Inches(11), Inches(0.7), size=36, bold=True)
rect(s, Inches(0.6), Inches(1.6), Inches(3.5), Inches(0.03), BRG)

sources = [
    ("Tavily API",    "Deep-search news retrieval, advanced relevancy ranking, real-time context for any theme"),
    ("NewsAPI",       "Broad news coverage, 7-day rolling window, sorted by relevancy, deduplicated against Tavily"),
    ("yfinance",      "OHLCV price data, 30/90/180-day lookback, volatility, returns, and average volume per ticker"),
    ("Anthropic SDK", "Direct Claude API calls via the official anthropic Python SDK for both Actor and Critic roles"),
    ("Dedalus Labs",  "OpenAI-compatible proxy endpoint for GPT models, initialised via openai SDK with custom base_url"),
]

y = Inches(1.95)
for name, desc in sources:
    rect(s, Inches(0.6), y + Inches(0.1), Inches(0.1), Inches(0.28), BRG)
    text(s, name, Inches(0.85), y, Inches(2.5), Inches(0.5), size=16, bold=True, color=GREEN_LT)
    text(s, desc, Inches(3.5), y, Inches(9.35), Inches(0.5), size=15, color=MUTED)
    y += Inches(0.9)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Tech Stack
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, BLACK)
green_bar(s)
text(s, "Technology Stack", Inches(0.6), Inches(0.8), Inches(11), Inches(0.7),
     size=36, bold=True)
rect(s, Inches(0.6), Inches(1.6), Inches(4.0), Inches(0.03), BRG)

cats = [
    ("Frontend",     ["Streamlit (Python)", "Custom CSS dark BRG theme", "Plotly (equity curves)"]),
    ("AI / LLMs",    ["Claude via Anthropic SDK", "GPT via OpenAI + Dedalus", "Adversarial multi-agent design"]),
    ("Data",         ["Tavily API", "NewsAPI", "yfinance OHLCV"]),
    ("Architecture", ["Actor-Critic pattern", "JSON-strict outputs", "Session-state pipeline"]),
]

xs = [Inches(0.5), Inches(3.8), Inches(7.1), Inches(10.4)]
for i, (cat, items) in enumerate(cats):
    x = xs[i]
    rect(s, x, Inches(2.0), Inches(2.85), Inches(4.8), SURFACE)
    rect(s, x, Inches(2.0), Inches(2.85), Inches(0.07), BRG)
    text(s, cat, x + Inches(0.15), Inches(2.15), Inches(2.6), Inches(0.5),
         size=16, bold=True, color=GREEN_LT)
    bullets(s, items, x + Inches(0.15), Inches(2.75), Inches(2.6), Inches(3.8), size=15)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — User Experience
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, BLACK)
green_bar(s)
text(s, "User Experience", Inches(0.6), Inches(0.8), Inches(11), Inches(0.7),
     size=36, bold=True)
rect(s, Inches(0.6), Inches(1.6), Inches(4.0), Inches(0.03), BRG)

ux_steps = [
    ("Step 1", "Input",              "Choose a theme (AI, DeFi, custom...), set risk appetite and investment horizon"),
    ("Step 2", "Market Intelligence","Live sentiment badge, asset snapshot table, and source headlines appear"),
    ("Step 3", "Agent Arena",        "Claude and GPT actors propose portfolios, critics challenge each other"),
    ("Step 4", "Backtest & Results", "Equity curve, win rate, max drawdown, and written AI performance insight"),
]

y = Inches(2.05)
for step, title, desc in ux_steps:
    rect(s, Inches(0.6), y, Inches(1.1), Inches(0.62), BRG)
    text(s, step, Inches(0.6), y, Inches(1.1), Inches(0.62),
         size=12, bold=True, align=PP_ALIGN.CENTER)
    text(s, title, Inches(1.85), y, Inches(2.6), Inches(0.62),
         size=18, bold=True, color=GREEN_LT)
    text(s, desc, Inches(4.6), y, Inches(8.2), Inches(0.62), size=16)
    if y < Inches(5.5):
        rect(s, Inches(1.12), y + Inches(0.62), Inches(0.06), Inches(0.48), BRG)
    y += Inches(1.12)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Roadmap
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, BLACK)
green_bar(s)
text(s, "Roadmap", Inches(0.6), Inches(0.8), Inches(11), Inches(0.7), size=36, bold=True)
rect(s, Inches(0.6), Inches(1.6), Inches(2.5), Inches(0.03), BRG)

phases = [
    ("Phase 1  \u2713", "Data Pipeline",
     ["Tavily + NewsAPI ingestion", "yfinance price data", "Sentiment synthesis"]),
    ("Phase 2  \u2713", "Frontend",
     ["Dark BRG Streamlit UI", "Single-page scroll layout", "Custom theme input"]),
    ("Phase 3  \u2192", "Agent Arena",
     ["Claude + GPT Actors", "Cross-model critic layer", "Risk Analyst synthesis"]),
    ("Phase 4  \u2192", "Backtesting",
     ["Paper-trading simulation", "Equity curve (Plotly)", "Sharpe / drawdown metrics"]),
]

xs2 = [Inches(0.5), Inches(3.8), Inches(7.1), Inches(10.4)]
for i, (phase, title, items) in enumerate(phases):
    x = xs2[i]
    done = "\u2713" in phase
    bc = BRG if done else SURFACE
    rect(s, x, Inches(2.0), Inches(2.85), Inches(5.1), SURFACE)
    rect(s, x, Inches(2.0), Inches(2.85), Inches(0.07), bc)
    text(s, phase, x + Inches(0.15), Inches(2.1), Inches(2.6), Inches(0.4),
         size=12, bold=True, color=BRG_LT if done else MUTED)
    text(s, title, x + Inches(0.15), Inches(2.6), Inches(2.6), Inches(0.5),
         size=16, bold=True, color=GREEN_LT if done else WHITE)
    bullets(s, items, x + Inches(0.15), Inches(3.2), Inches(2.6), Inches(3.7),
            size=14, color=WHITE if done else MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Closing
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, BLACK)
rect(s, 0, 0, Inches(0.35), H, BRG)
rect(s, Inches(0.35), Inches(3.15), W, Inches(0.04), BRG)

text(s, "KAPPA", Inches(0.65), Inches(1.2), Inches(11), Inches(1.8),
     size=80, bold=True)
text(s, "Built for the next generation of algorithmic traders.",
     Inches(0.65), Inches(3.3), Inches(11), Inches(0.65),
     size=22, italic=True, color=GREEN_LT)
text(s, "Actor-Critic  \u00b7  Cross-Model  \u00b7  Real-Time  \u00b7  Risk-Adjusted",
     Inches(0.65), Inches(4.15), Inches(11), Inches(0.55),
     size=16, color=MUTED)

prs.save("Kappa_Project_Deck.pptx")
print("Saved: Kappa_Project_Deck.pptx")
